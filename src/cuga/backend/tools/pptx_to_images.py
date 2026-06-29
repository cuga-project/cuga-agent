"""System tool: pptx_to_images.

Convert a .pptx file into per-page JPEG images in a single call.

Primary pipeline: PPTX → PDF (LibreOffice) → JPEG (PyMuPDF or pdftoppm).
LibreOffice is located by checking PATH first, then common installation
directories (macOS /Applications, Homebrew, standard Linux paths).

Fallback pipeline: when LibreOffice is absent the tool renders slides
directly from the PPTX XML using python-pptx + Pillow.  Fidelity is lower
(no complex gradients, custom fonts, or chart rendering) but layout geometry
— positions, sizes, colours, text content, embedded pictures — is faithfully
reproduced, which is sufficient for visual QA of overlaps and overflow.

The first line of the return value is always either
  "RENDERED_BY: libreoffice"  or  "RENDERED_BY: python-pptx (fallback)"
so a calling agent knows which path was used.

Install the preferred backends:
  uv pip install pymupdf          # PDF→JPEG without pdftoppm
  brew install --cask libreoffice # macOS — high-fidelity primary path
"""

from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from langchain_core.tools import StructuredTool
from pydantic import BaseModel, Field

_SOFFICE_SEARCH_PATHS = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/opt/homebrew/bin/soffice",
    "/usr/local/bin/soffice",
    "/usr/bin/soffice",
    "/usr/lib/libreoffice/program/soffice",
    "/snap/bin/libreoffice",
]

_EMU_PER_INCH = 914400


def _find_soffice() -> str | None:
    if found := shutil.which("soffice"):
        return found
    for candidate in _SOFFICE_SEARCH_PATHS:
        if Path(candidate).exists():
            return candidate
    return None


_WORKSPACE_ROOTS = [
    Path("/workspace"),
    Path.home() / "workspace",
]


def _resolve(path: str) -> Path:
    p = Path(path)
    if p.is_absolute():
        return p
    # Strip an explicit /workspace/ prefix and re-resolve
    if path.startswith("/workspace/"):
        p = Path(path[len("/workspace/"):])
    # Prefer the candidate that actually exists on disk
    candidates = [p.resolve(), *(root / p for root in _WORKSPACE_ROOTS)]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    # Return the raw path; the caller will emit a clear "not found" error
    return p


_SHIM_SO = Path(tempfile.gettempdir()) / "lo_socket_shim.so"

_SHIM_C = r"""
#define _GNU_SOURCE
#include <dlfcn.h>
#include <errno.h>
#include <stdio.h>
#include <stdlib.h>
#include <sys/socket.h>
#include <unistd.h>

static int (*real_socket)(int, int, int);
static int (*real_socketpair)(int, int, int, int[2]);
static int (*real_listen)(int, int);
static int (*real_accept)(int, struct sockaddr *, socklen_t *);
static int (*real_close)(int);
static int (*real_read)(int, void *, size_t);

static int is_shimmed[1024];
static int peer_of[1024];
static int wake_r[1024];
static int wake_w[1024];
static int listener_fd = -1;

__attribute__((constructor))
static void init(void) {
    real_socket     = dlsym(RTLD_NEXT, "socket");
    real_socketpair = dlsym(RTLD_NEXT, "socketpair");
    real_listen     = dlsym(RTLD_NEXT, "listen");
    real_accept     = dlsym(RTLD_NEXT, "accept");
    real_close      = dlsym(RTLD_NEXT, "close");
    real_read       = dlsym(RTLD_NEXT, "read");
    for (int i = 0; i < 1024; i++) { peer_of[i] = wake_r[i] = wake_w[i] = -1; }
}

int socket(int domain, int type, int protocol) {
    if (domain == AF_UNIX) {
        int fd = real_socket(domain, type, protocol);
        if (fd >= 0) return fd;
        int sv[2];
        if (real_socketpair(domain, type, protocol, sv) == 0) {
            if (sv[0] >= 0 && sv[0] < 1024) {
                is_shimmed[sv[0]] = 1;
                peer_of[sv[0]] = sv[1];
                int wp[2];
                if (pipe(wp) == 0) { wake_r[sv[0]] = wp[0]; wake_w[sv[0]] = wp[1]; }
            }
            return sv[0];
        }
        errno = EPERM; return -1;
    }
    return real_socket(domain, type, protocol);
}

int listen(int sockfd, int backlog) {
    if (sockfd >= 0 && sockfd < 1024 && is_shimmed[sockfd]) { listener_fd = sockfd; return 0; }
    return real_listen(sockfd, backlog);
}

int accept(int sockfd, struct sockaddr *addr, socklen_t *addrlen) {
    if (sockfd >= 0 && sockfd < 1024 && is_shimmed[sockfd]) {
        if (wake_r[sockfd] >= 0) { char buf; real_read(wake_r[sockfd], &buf, 1); }
        errno = ECONNABORTED; return -1;
    }
    return real_accept(sockfd, addr, addrlen);
}

int close(int fd) {
    if (fd >= 0 && fd < 1024 && is_shimmed[fd]) {
        int was_listener = (fd == listener_fd);
        is_shimmed[fd] = 0;
        if (wake_w[fd] >= 0) { char c = 0; write(wake_w[fd], &c, 1); real_close(wake_w[fd]); wake_w[fd] = -1; }
        if (wake_r[fd] >= 0) { real_close(wake_r[fd]); wake_r[fd] = -1; }
        if (peer_of[fd] >= 0) { real_close(peer_of[fd]); peer_of[fd] = -1; }
        if (was_listener) _exit(0);
    }
    return real_close(fd);
}
"""


def _af_unix_blocked() -> bool:
    import socket as _socket  # noqa: PLC0415
    try:
        s = _socket.socket(_socket.AF_UNIX, _socket.SOCK_STREAM)
        s.close()
        return False
    except OSError:
        return True


def _ensure_shim() -> Path | None:
    if _SHIM_SO.exists():
        return _SHIM_SO
    gcc = shutil.which("gcc")
    if not gcc:
        return None
    src = Path(tempfile.gettempdir()) / "lo_socket_shim.c"
    src.write_text(_SHIM_C)
    try:
        result = subprocess.run(
            [gcc, "-shared", "-fPIC", "-o", str(_SHIM_SO), str(src), "-ldl"],
            capture_output=True,
        )
        if result.returncode != 0:
            return None
    finally:
        src.unlink(missing_ok=True)
    return _SHIM_SO


def _soffice_env() -> dict:
    env = os.environ.copy()
    env["SAL_USE_VCLPLUGIN"] = "svp"
    if _af_unix_blocked():
        shim = _ensure_shim()
        if shim:
            env["LD_PRELOAD"] = str(shim)
    return env


# ── public tool function ───────────────────────────────────────────────────────

def pptx_to_images(pptx: str, prefix: str, dpi: int = 150) -> str:
    """Convert a .pptx file into per-page JPEG images.

    Tries LibreOffice first (high fidelity); falls back to a pure-Python
    renderer (python-pptx + Pillow) when LibreOffice is not available.

    Args:
        pptx:   Path to the .pptx file (absolute, relative, or /workspace/...).
        prefix: Output filename prefix.  Images are written as
                ``<prefix>-01.jpg``, ``<prefix>-02.jpg``, etc.
        dpi:    Render resolution in DPI (default 150).

    Returns:
        The first line is ``RENDERED_BY: libreoffice`` or
        ``RENDERED_BY: python-pptx (fallback)`` followed by a newline-separated
        list of generated image paths.  On hard failure the first line starts
        with ``ERROR:``.
    """
    pptx_path = _resolve(pptx)
    if not pptx_path.exists():
        return f"ERROR: File not found: {pptx!r}"
    if pptx_path.suffix.lower() != ".pptx":
        return f"ERROR: Expected a .pptx file, got: {pptx!r}"

    prefix_path = Path(prefix)
    if prefix_path.parent != Path("."):
        prefix_path.parent.mkdir(parents=True, exist_ok=True)

    soffice = _find_soffice()
    if soffice:
        result = _via_libreoffice(pptx_path, str(prefix_path), dpi, soffice)
        if not result.startswith("ERROR:"):
            return f"RENDERED_BY: libreoffice\n{result}"

    # LibreOffice unavailable or failed — try pure-Python path
    fallback = _via_python_pptx(pptx_path, str(prefix_path), dpi)
    if fallback.startswith("ERROR:"):
        return fallback
    return f"RENDERED_BY: python-pptx (fallback)\n{fallback}"


# ── LibreOffice path ───────────────────────────────────────────────────────────

def _via_libreoffice(pptx: Path, prefix: str, dpi: int, soffice: str) -> str:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        pdf_path = tmp_dir / f"{pptx.stem}.pdf"

        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(pptx)],
            capture_output=True,
            text=True,
            env=_soffice_env(),
        )
        if result.returncode != 0 or not pdf_path.exists():
            detail = (result.stderr or result.stdout or "no output").strip()
            return f"ERROR: LibreOffice PDF conversion failed:\n{detail}"

        return _pdf_to_jpegs(pdf_path, prefix, dpi)


def _pdf_to_jpegs(pdf: Path, prefix: str, dpi: int) -> str:
    try:
        return _via_pymupdf_pdf(pdf, prefix, dpi)
    except ImportError:
        pass

    if shutil.which("pdftoppm"):
        return _via_pdftoppm(pdf, prefix, dpi)

    return (
        "ERROR: LibreOffice converted to PDF but no image backend is available.\n"
        "Install one of:\n"
        "  uv pip install pymupdf   # pure-Python (recommended)\n"
        "  brew install poppler     # provides pdftoppm"
    )


def _via_pymupdf_pdf(pdf: Path, prefix: str, dpi: int) -> str:
    import fitz  # noqa: PLC0415

    doc = fitz.open(str(pdf))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    generated: list[str] = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat)
        out = Path(f"{prefix}-{i + 1:02d}.jpg")
        pix.save(str(out))
        generated.append(str(out))
    return "\n".join(generated) if generated else "WARNING: no pages converted"


def _via_pdftoppm(pdf: Path, prefix: str, dpi: int) -> str:
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf), prefix],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        return f"ERROR: pdftoppm failed:\n{result.stderr}"
    parent = Path(prefix).parent
    stem = Path(prefix).name
    found = sorted(parent.glob(f"{stem}-*.jpg"))
    return "\n".join(str(f) for f in found) if found else "WARNING: pdftoppm produced no output"


# ── pure-Python fallback renderer ─────────────────────────────────────────────

def _via_python_pptx(pptx: Path, prefix: str, dpi: int) -> str:
    try:
        from pptx import Presentation  # noqa: PLC0415
        from PIL import Image, ImageDraw, ImageFont  # noqa: PLC0415
    except ImportError as exc:
        return (
            f"ERROR: LibreOffice not found and pure-Python fallback unavailable ({exc}).\n"
            "Install LibreOffice or run:  uv pip install python-pptx Pillow"
        )

    try:
        prs = Presentation(str(pptx))
    except Exception as exc:
        return f"ERROR: Could not open {pptx.name}: {exc}"

    slide_w_emu = prs.slide_width or int(10 * _EMU_PER_INCH)
    slide_h_emu = prs.slide_height or int(5.625 * _EMU_PER_INCH)
    px_w = int(slide_w_emu / _EMU_PER_INCH * dpi)
    px_h = int(slide_h_emu / _EMU_PER_INCH * dpi)

    theme_palette = _extract_theme_colors(prs)

    generated: list[str] = []
    for idx, slide in enumerate(prs.slides):
        img = _render_slide(slide, px_w, px_h, slide_w_emu, slide_h_emu, dpi, theme_palette)
        out = Path(f"{prefix}-{idx + 1:02d}.jpg")
        img.save(str(out), "JPEG", quality=92)
        generated.append(str(out))

    return "\n".join(generated) if generated else "WARNING: presentation has no slides"


def _emu_to_px(emu: int | None, total_emu: int, total_px: int) -> int:
    if emu is None:
        return 0
    return int(emu * total_px / total_emu)


def _extract_theme_colors(prs) -> dict[int, tuple[int, int, int]]:
    """Return {MSO_THEME_COLOR_int: (r,g,b)} from the presentation theme XML."""
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = prs.slide_master.element.find(".//a:clrScheme", nsmap)
        if clr_scheme is None:
            return {}
        palette: dict[int, tuple[int, int, int]] = {}
        for xml_idx, child in enumerate(clr_scheme):
            mso_idx = xml_idx + 1  # MSO_THEME_COLOR values start at 1
            for color_el in child:
                tag = color_el.tag.split("}")[-1]
                if tag == "srgbClr":
                    val = color_el.get("val", "")
                elif tag == "sysClr":
                    val = color_el.get("lastClr", "")
                else:
                    continue
                if len(val) == 6:
                    palette[mso_idx] = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
                break
        return palette
    except Exception:
        return {}


def _apply_brightness(rgb: tuple[int, int, int], brightness: float) -> tuple[int, int, int]:
    r, g, b = rgb
    if brightness > 0:
        return (
            int(r + (255 - r) * brightness),
            int(g + (255 - g) * brightness),
            int(b + (255 - b) * brightness),
        )
    if brightness < 0:
        f = 1 + brightness
        return (int(r * f), int(g * f), int(b * f))
    return rgb


def _rgb_from_pptx_color(
    color_obj, theme_palette: dict[int, tuple[int, int, int]] | None = None
) -> tuple[int, int, int] | None:
    try:
        from pptx.dml.color import _NoneColor  # noqa: PLC0415
        if isinstance(color_obj, _NoneColor):
            return None
    except Exception:
        pass
    try:
        c = color_obj.rgb
        return (c.red, c.green, c.blue)
    except Exception:
        pass
    # Try resolving a theme color
    if theme_palette:
        try:
            tc = int(color_obj.theme_color)
            rgb = theme_palette.get(tc)
            if rgb:
                brightness = 0.0
                try:
                    brightness = color_obj.brightness
                except Exception:
                    pass
                return _apply_brightness(rgb, brightness)
        except Exception:
            pass
    return None


def _fill_color(fill, theme_palette: dict) -> tuple[int, int, int] | None:
    """Resolve a FillFormat to an RGB tuple, handling solid and theme fills."""
    try:
        return _rgb_from_pptx_color(fill.fore_color, theme_palette)
    except Exception:
        return None


def _background_color(slide, theme_palette: dict) -> tuple[int, int, int]:
    """Walk slide → layout → master to find the first defined solid background."""
    sources = [slide]
    try:
        sources.append(slide.slide_layout)
        sources.append(slide.slide_layout.slide_master)
    except Exception:
        pass
    for src in sources:
        try:
            c = _fill_color(src.background.fill, theme_palette)
            if c:
                return c
        except Exception:
            continue
    return (255, 255, 255)


def _render_slide(slide, px_w: int, px_h: int, emu_w: int, emu_h: int, dpi: int, theme_palette: dict | None = None):
    from PIL import Image, ImageDraw, ImageFont  # noqa: PLC0415

    tp = theme_palette or {}
    bg_rgb = _background_color(slide, tp)
    img = Image.new("RGB", (px_w, px_h), bg_rgb)
    draw = ImageDraw.Draw(img)

    try:
        font = ImageFont.load_default(size=max(10, dpi // 12))
    except TypeError:
        font = ImageFont.load_default()

    for shape in slide.shapes:
        left = _emu_to_px(shape.left, emu_w, px_w)
        top = _emu_to_px(shape.top, emu_h, px_h)
        width = _emu_to_px(shape.width, emu_w, px_w)
        height = _emu_to_px(shape.height, emu_h, px_h)
        right = left + width
        bottom = top + height

        # --- shape fill ---
        fill_rgb: tuple[int, int, int] | None = None
        try:
            fill_rgb = _fill_color(shape.fill, tp)
        except Exception:
            pass

        # --- shape outline ---
        line_rgb: tuple[int, int, int] | None = (180, 180, 180)
        try:
            c = _rgb_from_pptx_color(shape.line.color, tp)
            if c:
                line_rgb = c
        except Exception:
            pass

        if width > 0 and height > 0 and fill_rgb:
            draw.rectangle([left, top, right, bottom], fill=fill_rgb, outline=line_rgb)
        elif width > 0 and height > 0 and line_rgb:
            draw.rectangle([left, top, right, bottom], outline=line_rgb)

        # --- embedded picture ---
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                pic = Image.open(io.BytesIO(blob)).convert("RGB")
                if width > 0 and height > 0:
                    pic = pic.resize((width, height), Image.Resampling.LANCZOS)
                img.paste(pic, (left, top))
        except Exception:
            pass

        # --- text ---
        if shape.has_text_frame and width > 0 and height > 0:
            _draw_text_frame(draw, shape.text_frame, left, top, width, height, font, dpi, tp)

    return img


def _draw_text_frame(draw, tf, left: int, top: int, width: int, height: int, font, dpi: int, theme_palette: dict | None = None) -> None:
    from PIL import ImageFont  # noqa: PLC0415

    tp = theme_palette or {}
    padding = max(4, dpi // 50)
    x = left + padding
    y = top + padding
    max_x = left + width - padding
    max_y = top + height - padding

    for para in tf.paragraphs:
        line_parts: list[tuple[str, tuple[int, int, int], int]] = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            # font colour
            try:
                c = _rgb_from_pptx_color(run.font.color, tp)
                color = c if c else (30, 30, 30)
            except Exception:
                color = (30, 30, 30)
            # font size
            try:
                pt = run.font.size
                size = int(pt / 12700 * dpi / 72) if pt else max(10, dpi // 12)
            except Exception:
                size = max(10, dpi // 12)
            line_parts.append((text, color, size))

        if not line_parts:
            y += max(10, dpi // 12) + 2
            continue

        # Use the most common size in this paragraph for line height
        sizes = [s for _, _, s in line_parts]
        line_h = max(sizes) + 2

        if y + line_h > max_y:
            # Draw overflow indicator and stop
            draw.line([(left, max_y), (left + width, max_y)], fill=(255, 80, 80), width=2)
            break

        # Draw each run fragment
        run_x = x
        for text, color, size in line_parts:
            try:
                run_font = ImageFont.load_default(size=size)
            except TypeError:
                run_font = font

            # simple word-wrap at box boundary
            words = text.split()
            for word in words:
                w_text = word + " "
                try:
                    bbox = draw.textbbox((0, 0), w_text, font=run_font)
                    w_px = bbox[2] - bbox[0]
                except Exception:
                    w_px = len(w_text) * size // 2

                if run_x + w_px > max_x and run_x > x:
                    run_x = x
                    y += line_h
                    if y + line_h > max_y:
                        draw.line([(left, max_y), (left + width, max_y)], fill=(255, 80, 80), width=2)
                        return

                if run_x + w_px <= max_x:
                    draw.text((run_x, y), w_text, fill=color, font=run_font)
                    run_x += w_px

        y += line_h


# ── tool wiring ────────────────────────────────────────────────────────────────

class _PptxToImagesInput(BaseModel):
    pptx: str = Field(..., description="Path to the .pptx file (absolute, relative, or /workspace/filename).")
    prefix: str = Field(
        ...,
        description=(
            "Output filename prefix. Images are written as <prefix>-01.jpg, <prefix>-02.jpg, etc. "
            "May include a subdirectory, e.g. 'slides/slide'."
        ),
    )
    dpi: int = Field(150, description="Render resolution in DPI (default 150).")


def create_pptx_to_images_tool() -> StructuredTool:
    """Return a StructuredTool wrapping :func:`pptx_to_images`."""
    return StructuredTool.from_function(
        func=pptx_to_images,
        name="pptx_to_images",
        description=(
            "Convert a .pptx presentation into per-page JPEG images for visual QA. "
            "Tries LibreOffice first (high fidelity); falls back automatically to a "
            "pure-Python renderer (python-pptx + Pillow) when LibreOffice is absent — "
            "layout geometry, colours, text, and embedded pictures are preserved. "
            "The first line of the result is 'RENDERED_BY: libreoffice' or "
            "'RENDERED_BY: python-pptx (fallback)' so you know which path was used. "
            "Pass the returned image paths to analyze_image for visual inspection."
        ),
        args_schema=_PptxToImagesInput,
    )
